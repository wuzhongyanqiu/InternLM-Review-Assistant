import os
import logging
from langchain_community.document_loaders import UnstructuredMarkdownLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import fitz
import re
from docx import Document as dc
from langchain.schema import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor
import markdown2

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class SimpleParse:
    def __init__(self):
        self.chunk_size = 512
        self.overlap_size = 128
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.overlap_size
        )

    def clean_markdown(self, content):
        # 去掉不必要的内容，例如处理指令或特殊标签
        return re.sub(r'<.*?>', '', content)  # 去掉 HTML 标签

    def data_parsing(self, file_path):
        split_docs = []
        try:
            if file_path.lower().endswith(".md"):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                html_content = markdown2.markdown(content)  # 转换为 HTML
                text_content = BeautifulSoup(html_content, 'html.parser').get_text()
                text_content = self.clean_markdown(text_content.replace('\n\n', '\n'))
                split_docs = self.text_splitter.split_text(text_content)

            elif file_path.lower().endswith(".txt"):
                with open(file_path, 'r', encoding='utf-8') as file:
                    text_content = file.read()
                text_content = text_content.replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)

            elif file_path.lower().endswith(".docx"):
                doc = Document(file_path)
                text_content = '\n'.join(para.text for para in doc.paragraphs)
                text_content = text_content.replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)

            elif file_path.lower().endswith(".pptx"):
                extracted_text = []
                ppt = Presentation(file_path)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            extracted_text.append(shape.text)
                text_content = '\n'.join(extracted_text).replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)

            elif file_path.lower().endswith(".pdf"):
                pdf_document = fitz.open(file_path)
                text_content = ''
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    text_content += page.get_text()
                text_content = text_content.replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)

            else:
                logger.warning(f"Unsupported file type: {file_path}")
                return []

            # 在这里为每个分块添加文件名、路径和页数信息
            for idx, doc in enumerate(split_docs):
                split_docs[idx] = Document(
                    page_content=doc,
                    metadata={
                        "file_name": os.path.basename(file_path),
                        "file_path": file_path,
                        "page_num": idx + 1
                    }
                )

        except Exception as e:
            logger.error(f"An error occurred while processing the file {file_path}: {e}")
            raise

        return split_docs
        
    def process_folder(self, folder_path):
        aggregated_data = []
        file_paths = [os.path.join(folder_path, filename) for filename in os.listdir(folder_path)]
        with ThreadPoolExecutor() as executor:
            results = executor.map(self.data_parsing, file_paths)
            for result in results:
                aggregated_data.extend(result)
        logger.info("process_folder is finished!")
        return aggregated_data

if __name__ == "__main__":
    # Example usage
    pass
    # Process or save the aggregated_data as needed
