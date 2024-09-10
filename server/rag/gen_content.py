import os
import logging
from langchain_community.document_loaders import UnstructuredMarkdownLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor
from ..minerU_worker import MinerU

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class GenContentDB:
    def __init__(self):
        self.chunk_size = 256
        self.overlap_size = 50
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.overlap_size
        )

    def data_parsing(self, file_path):
        split_docs = []
        try:                
            if file_path.lower().endswith(".md"):
                loader = UnstructuredMarkdownLoader(file_path)
                md_pages = loader.load()
                md_page = md_pages[0]
                md_page.page_content = md_page.page_content.replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(md_page.page_content)

            elif file_path.lower().endswith(".txt"):
                with open(file_path, 'r', encoding='utf-8') as file:
                    text_content = file.read()
                text_content = text_content.replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)

            elif file_path.lower().endswith(".docx"):
                doc = Document(file_path)
                text_content = '\n'.join(para.text for para in doc.paragraphs)
                text_content = text_content.replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)

            elif file_path.lower().endswith((".html", ".htm")):
                with open(file_path, 'r', encoding='utf-8') as file:
                    html_content = file.read()
                soup = BeautifulSoup(html_content, 'lxml')
                text_content = soup.get_text(separator='\n')
                text_content = text_content.replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)

            elif file_path.lower().endswith(".pptx"):
                extracted_text = []
                ppt = Presentation(file_path)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            extracted_text.append(shape.text)
                text_content = '\n'.join(extracted_text).replace('\n\n', '\n')
                split_docs = self.text_splitter.split_text(text_content)
            else:
                logger.warning(f"Unsupported file type: {file_path}")
                
        except Exception as e:
            logger.error(f"An error occurred while processing the file {file_path}: {e}")
            raise

        return split_docs
        
    def process_folder(self, folder_path):
        aggregated_data = []
        file_paths = [os.path.join(folder_path, filename) for filename in os.listdir(folder_path)]
        with ThreadPoolExecutor() as executor:
            results = executor.map(self.data_parsing, file_paths)
            for result in results:
                aggregated_data.extend(result)
        logger.info("process_folder is finished!")
        return aggregated_data

if __name__ == "__main__":
    # Example usage
    pass
    # Process or save the aggregated_data as needed
